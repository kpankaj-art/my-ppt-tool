import io
import re
import time
import pandas as pd
import openpyxl
from openpyxl.styles import PatternFill
from pptx import Presentation
import streamlit as st

# Page Setup
st.set_page_config(page_title="Pro PPT & Excel Matcher", page_icon="⚡", layout="centered")

# --- PROFESSIONAL UI & CUSTOM CSS ---
st.markdown("""
    <style>
    [data-testid="stHeader"] { display: none !important; }
    #MainMenu { visibility: hidden; }
    header { visibility: hidden; }
    footer { visibility: hidden; }
    div[class*="stAppHeader"] { display: none !important; }

    .main { background-color: #0e1117; }
    
    .css-card {
        background: linear-gradient(135deg, #1e2638 0%, #111827 100%);
        border: 1px solid #374151;
        border-radius: 12px;
        padding: 24px;
        margin-bottom: 20px;
        box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.3);
    }
    
    .title-text {
        font-family: 'Inter', sans-serif;
        font-size: 28px;
        font-weight: 700;
        color: #f9fafb;
        margin-bottom: 8px;
    }
    .subtitle-text {
        font-size: 14px;
        color: #9ca3af;
        margin-bottom: 24px;
    }
    
    .metric-container {
        display: flex;
        justify-content: space-between;
        gap: 12px;
        margin-top: 15px;
        margin-bottom: 15px;
    }
    .metric-box {
        background: #1f2937;
        border-radius: 8px;
        padding: 12px 16px;
        text-align: center;
        flex: 1;
        border: 1px solid #374151;
    }
    .metric-value {
        font-size: 20px;
        font-weight: 700;
    }
    .metric-label {
        font-size: 12px;
        color: #9ca3af;
        margin-top: 4px;
    }
    
    .stDownloadButton > button {
        width: 100%;
        border-radius: 8px;
        font-weight: 600;
        height: 48px;
    }
    </style>
""", unsafe_allow_html=True)

# --- HELPER FUNCTIONS ---

def normalize_str(val):
    """Remove spaces, hyphens, and extra symbols for accurate matching"""
    if pd.isna(val) or str(val).strip().lower() in ['nan', 'none', '']:
        return ""
    val_str = str(val).strip().lower()
    # Remove all whitespace characters
    val_str = re.sub(r'\s+', '', val_str)
    # Remove decimal trailing .0 if present
    if val_str.endswith('.0'):
        val_str = val_str[:-2]
    return val_str

def load_any_excel_file(uploaded_file):
    """Universal File Reader for .xlsx, .xls, .xlsm, .csv"""
    filename = uploaded_file.name.lower()
    file_bytes = uploaded_file.getvalue()
    
    if filename.endswith('.csv'):
        df = pd.read_csv(io.BytesIO(file_bytes), dtype=str)
    else:
        try:
            df = pd.read_excel(io.BytesIO(file_bytes), dtype=str)
        except Exception:
            try:
                df = pd.read_excel(io.BytesIO(file_bytes), engine='openpyxl', dtype=str)
            except Exception:
                df = pd.read_excel(io.BytesIO(file_bytes), engine='xlrd', dtype=str)
                
    wb = openpyxl.Workbook()
    ws = wb.active
    
    # Write Headers
    ws.append(list(df.columns))
    
    # Write Data Rows
    for row in df.itertuples(index=False):
        ws.append(list(row))
        
    return df, wb

def extract_text_from_slide(slide):
    text_runs = []
    def process_shape(shape):
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text:
                    text_runs.append(paragraph.text.strip())
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text:
                        text_runs.append(cell.text.strip())
        elif shape.shape_type == 6:  # Group Shape
            for sub_shape in shape.shapes:
                process_shape(sub_shape)

    for shape in slide.shapes:
        process_shape(shape)
    return "\n".join(text_runs)

def extract_details_from_text(text):
    name = ""
    name_match = re.search(r'(?:Outlet\s*Name|Party\s*Name|Customer\s*Name|Dealer\s*Name|Shop\s*Name|Store\s*Name|Name)\s*[:\-]\s*([^\n\r]+)', text, re.IGNORECASE)
    
    if name_match:
        name = name_match.group(1).strip()
        name = re.split(r'(?:Address|Contact|City|Date|Qty|Size|Type|Width|Height)', name, flags=re.IGNORECASE)[0].strip()
    else:
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        for l in lines:
            if not any(k in l.lower() for k in ['qty', 'size', 'type', 'contact', 'address', 'city', 'date', 'width', 'height']):
                name = l[:40].strip()
                break

    mob_match = re.search(r'\b[6-9]\d{9}\b', text)
    mobile = mob_match.group(0) if mob_match else ""

    sap_match = re.search(r'\b\d{6,10}\b', text)
    sap = sap_match.group(0) if sap_match and sap_match.group(0) != mobile else ""

    dim_match = re.search(r'(\d+(?:\.\d+)?)\s*x\s*(\d+(?:\.\d+)?)', text, re.IGNORECASE)
    w = dim_match.group(1) if dim_match else ""
    h = dim_match.group(2) if dim_match else ""

    return name, mobile, sap, w, h

def process_files_with_progress(pptx_file, excel_file, progress_bar, status_text):
    status_text.markdown("⏳ **Reading Excel Data...**")
    progress_bar.progress(10)
    
    df, wb = load_any_excel_file(excel_file)
    
    # Column Identification
    col_mobile = next((c for c in df.columns if any(k in str(c).lower() for k in ['mobile', 'contact', 'phone', 'num'])), None)
    col_sap = next((c for c in df.columns if any(k in str(c).lower() for k in ['sap', 'code', 'dealer_code', 'id'])), None)
    col_name = next((c for c in df.columns if any(k in str(c).lower() for k in ['outlet', 'party', 'customer', 'dealer', 'shop', 'store', 'name'])), None)
    col_city = next((c for c in df.columns if any(k in str(c).lower() for k in ['city', 'location', 'address', 'town', 'place'])), None)
    col_width = next((c for c in df.columns if 'width' in str(c).lower() or str(c).strip().lower() in ['w', 'w(ft)', 'width(ft)']), None)
    col_height = next((c for c in df.columns if 'height' in str(c).lower() or str(c).strip().lower() in ['h', 'h(ft)', 'height(ft)']), None)

    excel_criteria = []
    for idx, row in df.iterrows():
        raw_mob = str(row[col_mobile]) if col_mobile and pd.notna(row[col_mobile]) else ""
        raw_sap = str(row[col_sap]) if col_sap and pd.notna(row[col_sap]) else ""
        raw_name = str(row[col_name]) if col_name and pd.notna(row[col_name]) else ""
        raw_city = str(row[col_city]) if col_city and pd.notna(row[col_city]) else ""
        raw_w = row[col_width] if col_width else ""
        raw_h = row[col_height] if col_height else ""

        excel_criteria.append({
            'row_idx': idx, 
            'mobile': normalize_str(raw_mob), 
            'sap': normalize_str(raw_sap), 
            'name': normalize_str(raw_name), 
            'city': normalize_str(raw_city), 
            'width': normalize_str(raw_w), 
            'height': normalize_str(raw_h)
        })

    status_text.markdown("⏳ **Analyzing Presentation Slides...**")
    progress_bar.progress(25)
    prs = Presentation(pptx_file)
    slides = list(prs.slides)
    total_slides = len(slides)
    
    raw_slide_texts = []
    for idx, slide in enumerate(slides):
        raw_slide_texts.append(extract_text_from_slide(slide))
        curr_p = 25 + int((idx + 1) / total_slides * 20)
        progress_bar.progress(min(curr_p, 45))
        
    # Standard Normalized Slide Texts (Spaces Removed for Space-Independent Matching)
    slide_texts_norm = [normalize_str(t) for t in raw_slide_texts]

    status_text.markdown("⏳ **Smart Matching Slides (Space Insensitive)...**")
    matched_indices = []
    seen_slides = set()
    matched_excel_rows = set()

    total_items = len(excel_criteria)
    for i, item in enumerate(excel_criteria):
        row_idx = item['row_idx']
        mob, sap = item['mobile'], item['sap']
        name, city = item['name'], item['city']
        w, h = item['width'], item['height']

        best_match_idx = None

        # 1. Primary: Match Unique Identifier (Mobile / SAP) + Size (Width & Height)
        if w and h:
            for idx, norm_text in enumerate(slide_texts_norm):
                if idx in seen_slides: continue
                size_match = (w in norm_text and h in norm_text) or (f"{w}x{h}" in norm_text)
                if size_match:
                    if (mob and mob in norm_text) or (sap and sap in norm_text):
                        best_match_idx = idx
                        break

        # 2. Match Unique Identifier (Mobile / SAP) Without Size
        if best_match_idx is None:
            if mob:
                for idx, norm_text in enumerate(slide_texts_norm):
                    if idx not in seen_slides and mob in norm_text:
                        best_match_idx = idx
                        break
            elif sap:
                for idx, norm_text in enumerate(slide_texts_norm):
                    if idx not in seen_slides and sap in norm_text:
                        best_match_idx = idx
                        break

        # 3. CRITICAL FALLBACK: Name + Size Match (When Mobile & SAP are Blank)
        if best_match_idx is None and name and len(name) > 2 and w and h:
            for idx, norm_text in enumerate(slide_texts_norm):
                if idx in seen_slides: continue
                size_match = (w in norm_text and h in norm_text) or (f"{w}x{h}" in norm_text)
                if name in norm_text and size_match:
                    best_match_idx = idx
                    break

        # 4. Fallback: Name + City
        if best_match_idx is None and name and len(name) > 2 and city:
            for idx, norm_text in enumerate(slide_texts_norm):
                if idx not in seen_slides and name in norm_text and city in norm_text:
                    best_match_idx = idx
                    break

        # 5. Last Resort: Name Match Only
        if best_match_idx is None and name and len(name) > 2:
            for idx, norm_text in enumerate(slide_texts_norm):
                if idx not in seen_slides and name in norm_text:
                    best_match_idx = idx
                    break

        if best_match_idx is not None:
            matched_indices.append(best_match_idx)
            seen_slides.add(best_match_idx)
            matched_excel_rows.add(row_idx)

        curr_p = 45 + int((i + 1) / total_items * 30)
        progress_bar.progress(min(curr_p, 75))

    status_text.markdown("⏳ **Reordering PPT Slides...**")
    progress_bar.progress(80)
    
    sldIdLst = prs.slides._sldIdLst
    original_sldIds = list(sldIdLst)
    for sldId in original_sldIds:
        sldIdLst.remove(sldId)

    for idx in matched_indices:
        sldIdLst.append(original_sldIds[idx])

    unmatched_slide_indices = []
    for idx, sldId in enumerate(original_sldIds):
        if idx not in matched_indices:
            sldIdLst.append(sldId)
            unmatched_slide_indices.append(idx)

    out_pptx_io = io.BytesIO()
    prs.save(out_pptx_io)
    out_pptx_io.seek(0)

    status_text.markdown("⏳ **Generating Final Excel Report...**")
    progress_bar.progress(90)
    
    ws = wb.active
    headers = [cell.value for cell in ws[1]]
    remark_col_num = len(headers) + 1
    ws.cell(row=1, column=remark_col_num, value="Remark")

    red_fill = PatternFill(start_color="FF9999", end_color="FF9999", fill_type="solid")
    green_fill = PatternFill(start_color="99FF99", end_color="99FF99", fill_type="solid")

    missing_count = 0
    for idx, row in enumerate(df.iterrows()):
        excel_row_num = idx + 2
        if idx not in matched_excel_rows:
            missing_count += 1
            for col_num in range(1, remark_col_num + 1):
                ws.cell(row=excel_row_num, column=col_num).fill = red_fill
            ws.cell(row=excel_row_num, column=remark_col_num, value="Not Found in PPT")

    extra_count = 0
    for s_idx in unmatched_slide_indices:
        extra_count += 1
        slide_txt = raw_slide_texts[s_idx]
        ext_name, ext_mob, ext_sap, ext_w, ext_h = extract_details_from_text(slide_txt)

        new_row_idx = ws.max_row + 1

        if col_name:
            c_idx = df.columns.get_loc(col_name) + 1
            ws.cell(row=new_row_idx, column=c_idx, value=ext_name)
        elif len(headers) >= 1:
            ws.cell(row=new_row_idx, column=1, value=ext_name)

        if col_mobile:
            c_idx = df.columns.get_loc(col_mobile) + 1
            ws.cell(row=new_row_idx, column=c_idx, value=ext_mob)

        if col_sap:
            c_idx = df.columns.get_loc(col_sap) + 1
            ws.cell(row=new_row_idx, column=c_idx, value=ext_sap)

        if col_width:
            c_idx = df.columns.get_loc(col_width) + 1
            ws.cell(row=new_row_idx, column=c_idx, value=ext_w)

        if col_height:
            c_idx = df.columns.get_loc(col_height) + 1
            ws.cell(row=new_row_idx, column=c_idx, value=ext_h)

        ws.cell(row=new_row_idx, column=remark_col_num, value="ye ppt me extra hai")

        for col_num in range(1, remark_col_num + 1):
            ws.cell(row=new_row_idx, column=col_num).fill = green_fill

    out_excel_io = io.BytesIO()
    wb.save(out_excel_io)
    out_excel_io.seek(0)

    progress_bar.progress(100)
    status_text.markdown("✅ **Processing Complete!**")
    time.sleep(0.5)

    return out_pptx_io.getvalue(), out_excel_io.getvalue(), len(matched_indices), missing_count, extra_count

# --- MAIN INTERFACE ---

st.markdown("""
    <div class="css-card">
        <div class="title-text">⚡ Smart PPT & Excel Matcher Pro</div>
        <div class="subtitle-text">Match and reorder your presentation slides with Excel data seamlessly.</div>
    </div>
""", unsafe_allow_html=True)

uploaded_pptx = st.file_uploader("1. Upload PowerPoint Presentation (.pptx)", type=["pptx"])
uploaded_excel = st.file_uploader("2. Upload Master Excel File (.xlsx, .xls, .csv, .xlsm)", type=["xlsx", "xls", "csv", "xlsm"])

st.markdown("<br>", unsafe_allow_html=True)

if st.button("🚀 Process & Sync Files", type="primary", use_container_width=True):
    if uploaded_pptx and uploaded_excel:
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        try:
            out_pptx_bytes, out_excel_bytes, matched_cnt, missing_cnt, extra_cnt = process_files_with_progress(
                uploaded_pptx, uploaded_excel, progress_bar, status_text
            )
            
            st.session_state["out_pptx"] = out_pptx_bytes
            st.session_state["out_excel"] = out_excel_bytes
            st.session_state["matched_cnt"] = matched_cnt
            st.session_state["missing_cnt"] = missing_cnt
            st.session_state["extra_cnt"] = extra_cnt
            st.session_state["processed"] = True
            
            status_text.empty()
            progress_bar.empty()
            
        except Exception as e:
            status_text.empty()
            progress_bar.empty()
            st.error(f"❌ Error during processing: {str(e)}")
    else:
        st.warning("⚠️ Kripya dono PowerPoint aur Excel files upload karein!")

# --- DISPLAY RESULTS & DOWNLOADS ---
if st.session_state.get("processed", False):
    st.markdown(f"""
        <div class="metric-container">
            <div class="metric-box">
                <div class="metric-value" style="color: #34d399;">{st.session_state['matched_cnt']}</div>
                <div class="metric-label">Matched Slides</div>
            </div>
            <div class="metric-box">
                <div class="metric-value" style="color: #f87171;">{st.session_state['missing_cnt']}</div>
                <div class="metric-label">Missing (Red)</div>
            </div>
            <div class="metric-box">
                <div class="metric-value" style="color: #60a5fa;">{st.session_state['extra_cnt']}</div>
                <div class="metric-label">Extra PPT (Green)</div>
            </div>
        </div>
    """, unsafe_allow_html=True)

    col1, col2 = st.columns(2)
    with col1:
        st.download_button(
            label="📥 Download Sorted PPT",
            data=st.session_state["out_pptx"],
            file_name="Sorted_Presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="btn_ppt"
        )
    with col2:
        st.download_button(
            label="📥 Download Updated Excel",
            data=st.session_state["out_excel"],
            file_name="Missing_Report.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key="btn_excel"
        )
